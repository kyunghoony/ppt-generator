from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches
from typing import List, Optional
from .utils import hex_to_rgb


def add_chart_to_slide(slide, chart_type: str, data: dict, labels: list,
                       left, top, width, height,
                       chart_colors: Optional[List[str]] = None):
    """Adds a chart to the given slide with optional preset colors."""
    chart_data = CategoryChartData()
    chart_data.categories = labels

    for series_name, series_values in data.items():
        chart_data.add_series(series_name, series_values)

    ctype = XL_CHART_TYPE.COLUMN_CLUSTERED
    if chart_type == "line":
        ctype = XL_CHART_TYPE.LINE
    elif chart_type == "pie":
        ctype = XL_CHART_TYPE.PIE

    chart_shape = slide.shapes.add_chart(
        ctype, left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    # Apply preset colors to series
    if chart_colors:
        for i, series in enumerate(chart.series):
            if i < len(chart_colors):
                color = hex_to_rgb(chart_colors[i])
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
                if hasattr(series.format, 'line'):
                    series.format.line.color.rgb = color

    return chart
