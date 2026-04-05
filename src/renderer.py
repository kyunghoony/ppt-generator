from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from typing import List, Dict, Any
from .primitives import (
    SlidePrimitive, TitleSlide, SectionDivider, ContentSlide,
    TwoColumnSlide, ImageSlide, ChartSlide, TableSlide,
    ComparisonSlide, TimelineSlide, MetricSlide, BlankSlide
)
from .utils import hex_to_rgb, apply_text_style
from .charts import add_chart_to_slide


class Renderer:
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.prs = Presentation()

        # Set slide size
        width = config.get("slide_size", {}).get("width", 13.333)
        height = config.get("slide_size", {}).get("height", 7.5)
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)
        self.slide_w = width
        self.slide_h = height

        # Layouts
        self.blank_layout = self.prs.slide_layouts[6]

        # Config shortcuts
        self.colors = config.get("colors", {})
        self.fonts = config.get("fonts", {})
        self.spacing = config.get("spacing", {})
        self.footer_cfg = config.get("footer", {})

        # Track slide index for footer page numbers
        self._slide_index = 0

    def render(self, primitives: List[SlidePrimitive], output_path: str):
        for primitive in primitives:
            self._render_slide(primitive)
        self.prs.save(output_path)

    def _render_slide(self, primitive: SlidePrimitive):
        self._slide_index += 1

        if isinstance(primitive, TitleSlide):
            self._render_title(primitive)
        elif isinstance(primitive, SectionDivider):
            self._render_section_divider(primitive)
        elif isinstance(primitive, ContentSlide):
            self._render_content(primitive)
        elif isinstance(primitive, ChartSlide):
            self._render_chart(primitive)
        elif isinstance(primitive, MetricSlide):
            self._render_metrics(primitive)
        elif isinstance(primitive, ComparisonSlide):
            self._render_comparison(primitive)
        elif isinstance(primitive, TwoColumnSlide):
            self._render_two_column(primitive)
        elif isinstance(primitive, TableSlide):
            self._render_table(primitive)
        elif isinstance(primitive, TimelineSlide):
            self._render_timeline(primitive)
        elif isinstance(primitive, BlankSlide):
            self._render_blank(primitive)
        else:
            self._render_fallback(primitive)

    # ── helpers ─────────────────────────────────────────────────

    def _set_slide_bg(self, slide, color_hex: str):
        """Set solid background color on a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color_hex)

    def _add_rounded_rect(self, slide, left, top, width, height, fill_hex, corner_radius=Inches(0.15)):
        """Add a rounded rectangle shape with fill color."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        shape.line.fill.background()  # No border
        # Set corner radius via XML
        sp = shape._element
        prstGeom = sp.find(qn('a:prstGeom'), sp.nsmap) if hasattr(sp, 'nsmap') else None
        if prstGeom is None:
            for child in sp.iter():
                if child.tag.endswith('prstGeom'):
                    prstGeom = child
                    break
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = prstGeom.makeelement(qn('a:avLst'), {})
                prstGeom.append(avLst)
            # Clear existing adjustments
            for gd in list(avLst):
                avLst.remove(gd)
            # Corner radius as percentage (50000 = max round, lower = less round)
            gd = avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': 'val 10000'})
            avLst.append(gd)
        return shape

    def _add_textbox(self, slide, left, top, width, height):
        """Add a textbox and return it."""
        return slide.shapes.add_textbox(left, top, width, height)

    def _add_footer(self, slide):
        """Add footer to slide (left: text, right: page number)."""
        if not self.footer_cfg.get("enabled", False):
            return

        footer_text = self.footer_cfg.get("text", "")
        show_page = self.footer_cfg.get("show_page_number", True)
        color = self.colors.get("text_secondary", "#555555")
        font_name = self.fonts.get("body", "Arial")
        y = Inches(self.slide_h - 0.3)
        margin_l = Inches(self.spacing.get("margin_left", 0.8))
        margin_r = Inches(self.spacing.get("margin_right", 0.8))
        content_w = Inches(self.slide_w) - margin_l - margin_r

        if footer_text:
            box = slide.shapes.add_textbox(margin_l, y, Inches(4), Inches(0.3))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = footer_text
            apply_text_style(run, font_name, 9, color)

        if show_page:
            box = slide.shapes.add_textbox(
                Inches(self.slide_w) - margin_r - Inches(1.5), y, Inches(1.5), Inches(0.3)
            )
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = str(self._slide_index)
            apply_text_style(run, font_name, 9, color)

    # ── Title Slide ─────────────────────────────────────────────

    def _render_title(self, prim: TitleSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        self._set_slide_bg(slide, primary)

        font_name = self.fonts.get("title", "Arial")
        text_light = self.colors.get("text_light", "#FFFFFF")
        text_secondary = self.colors.get("text_secondary", "#555555")

        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title — 44pt bold white, vertically centered slightly above
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(2.2), Inches(content_w), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = prim.title
        apply_text_style(run, font_name, 44, text_light)
        run.font.bold = True

        # Subtitle — 18pt white
        if prim.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(margin_l), Inches(4.4), Inches(content_w), Inches(1)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = prim.subtitle
            apply_text_style(run, font_name, 18, text_light)
        # No footer on title slide

    # ── Section Divider ─────────────────────────────────────────

    def _render_section_divider(self, prim: SectionDivider):
        slide = self.prs.slides.add_slide(self.blank_layout)
        secondary = self.colors.get("secondary", "#16213e")
        self._set_slide_bg(slide, secondary)

        font_name = self.fonts.get("title", "Arial")
        text_light = self.colors.get("text_light", "#FFFFFF")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title — 36pt bold white, vertically centered
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(2.5), Inches(content_w), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = prim.title
        apply_text_style(run, font_name, 36, text_light)
        run.font.bold = True

        # Subtitle — 16pt white
        if prim.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(margin_l), Inches(4.2), Inches(content_w), Inches(1)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = prim.subtitle
            apply_text_style(run, font_name, 16, text_light)
        # No footer on section divider

    # ── Content Slide ───────────────────────────────────────────

    def _render_content(self, prim: ContentSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        text_primary = self.colors.get("text_primary", "#1a1a2e")
        accent = self.colors.get("accent", "#0f3460")
        font_title = self.fonts.get("title", "Arial")
        font_body = self.fonts.get("body", "Arial")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title — 28pt bold primary
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(0.5), Inches(content_w), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = prim.title
        apply_text_style(run, font_title, 28, primary)
        run.font.bold = True

        # Body — 16pt
        body_top = 1.8
        body_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(body_top), Inches(content_w), Inches(4.5)
        )
        tf = body_box.text_frame
        tf.word_wrap = True

        if prim.body:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = prim.body
            apply_text_style(run, font_body, 16, text_primary)
            p.space_after = Pt(8)

        if prim.bullets:
            for bullet_text in prim.bullets:
                p = tf.add_paragraph()
                p.space_before = Pt(4)
                p.space_after = Pt(4)
                p.line_spacing = Pt(24)  # 1.5x of 16pt
                # Bullet marker in accent color
                bullet_run = p.add_run()
                bullet_run.text = "  \u2022  "
                apply_text_style(bullet_run, font_body, 16, accent)
                # Bullet text
                text_run = p.add_run()
                text_run.text = bullet_text
                apply_text_style(text_run, font_body, 16, text_primary)

        self._add_footer(slide)

    # ── Metrics Slide ───────────────────────────────────────────

    def _render_metrics(self, prim: MetricSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        secondary = self.colors.get("text_secondary", "#555555")
        card_bg = self.colors.get("card_background", "#F5F5F7")
        accent_pos = self.colors.get("accent_positive", "#16A34A")
        accent_neg = self.colors.get("accent_negative", "#DC2626")
        font_title = self.fonts.get("title", "Arial")
        font_body = self.fonts.get("body", "Arial")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title — 28pt bold primary
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(0.5), Inches(content_w), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.title
        apply_text_style(run, font_title, 28, primary)
        run.font.bold = True

        # Metric cards
        num = len(prim.metrics)
        if num == 0:
            self._add_footer(slide)
            return

        gap = 0.3  # gap between cards
        total_gap = gap * (num - 1)
        card_w = (content_w - total_gap) / num
        card_h = 3.0
        card_top = 2.2

        for i, metric in enumerate(prim.metrics):
            left = margin_l + i * (card_w + gap)

            # Card background
            self._add_rounded_rect(
                slide, Inches(left), Inches(card_top),
                Inches(card_w), Inches(card_h), card_bg
            )

            # Label — 12pt secondary
            label_box = slide.shapes.add_textbox(
                Inches(left + 0.3), Inches(card_top + 0.3),
                Inches(card_w - 0.6), Inches(0.4)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = metric.get("label", "")
            apply_text_style(run, font_body, 12, secondary)

            # Value — 40pt bold primary
            val_box = slide.shapes.add_textbox(
                Inches(left + 0.3), Inches(card_top + 0.8),
                Inches(card_w - 0.6), Inches(1.2)
            )
            tf = val_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = metric.get("value", "")
            apply_text_style(run, font_body, 40, primary)
            run.font.bold = True

            # Delta — 14pt colored
            delta = metric.get("delta", "")
            if delta:
                delta_box = slide.shapes.add_textbox(
                    Inches(left + 0.3), Inches(card_top + 2.1),
                    Inches(card_w - 0.6), Inches(0.5)
                )
                tf = delta_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                # Determine color based on sign
                if delta.startswith("-") or delta.startswith("\u2212"):
                    delta_color = accent_neg
                    run.text = delta
                else:
                    delta_color = accent_pos
                    run.text = delta
                apply_text_style(run, font_body, 14, delta_color)
                run.font.bold = True

        self._add_footer(slide)

    # ── Comparison Slide ────────────────────────────────────────

    def _render_comparison(self, prim: ComparisonSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        accent = self.colors.get("accent", "#0f3460")
        card_bg = self.colors.get("card_background", "#F5F5F7")
        text_light = self.colors.get("text_light", "#FFFFFF")
        font_title = self.fonts.get("title", "Arial")
        font_body = self.fonts.get("body", "Arial")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title — 28pt bold primary
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(0.5), Inches(content_w), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.title
        apply_text_style(run, font_title, 28, primary)
        run.font.bold = True

        # Two cards: 45% width each, 4% gap
        gap_pct = 0.04
        card_pct = 0.45
        card_w = content_w * card_pct
        gap_w = content_w * gap_pct
        # Center the pair
        total_w = card_w * 2 + gap_w
        start_x = margin_l + (content_w - total_w) / 2
        card_top = 1.8
        card_h = 4.5

        # Left card — accent background, white text
        self._add_rounded_rect(
            slide, Inches(start_x), Inches(card_top),
            Inches(card_w), Inches(card_h), accent
        )
        # Left label
        lbl_box = slide.shapes.add_textbox(
            Inches(start_x + 0.3), Inches(card_top + 0.3),
            Inches(card_w - 0.6), Inches(0.5)
        )
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.left_label
        apply_text_style(run, font_body, 16, text_light)
        run.font.bold = True

        # Left items
        items_box = slide.shapes.add_textbox(
            Inches(start_x + 0.3), Inches(card_top + 1.0),
            Inches(card_w - 0.6), Inches(card_h - 1.3)
        )
        tf = items_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(prim.left_items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"\u2022  {item}"
            apply_text_style(run, font_body, 14, text_light)

        # Right card — light gray background, primary text
        right_x = start_x + card_w + gap_w
        self._add_rounded_rect(
            slide, Inches(right_x), Inches(card_top),
            Inches(card_w), Inches(card_h), card_bg
        )
        # Right label
        lbl_box = slide.shapes.add_textbox(
            Inches(right_x + 0.3), Inches(card_top + 0.3),
            Inches(card_w - 0.6), Inches(0.5)
        )
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.right_label
        apply_text_style(run, font_body, 16, primary)
        run.font.bold = True

        # Right items
        items_box = slide.shapes.add_textbox(
            Inches(right_x + 0.3), Inches(card_top + 1.0),
            Inches(card_w - 0.6), Inches(card_h - 1.3)
        )
        tf = items_box.text_frame
        tf.word_wrap = True
        text_primary = self.colors.get("text_primary", "#1a1a2e")
        for j, item in enumerate(prim.right_items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"\u2022  {item}"
            apply_text_style(run, font_body, 14, text_primary)

        self._add_footer(slide)

    # ── Two Column Slide ────────────────────────────────────────

    def _render_two_column(self, prim: TwoColumnSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        text_primary = self.colors.get("text_primary", "#1a1a2e")
        divider_color = self.colors.get("divider", "#E5E7EB")
        font_title = self.fonts.get("title", "Arial")
        font_body = self.fonts.get("body", "Arial")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(0.5), Inches(content_w), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.title
        apply_text_style(run, font_title, 28, primary)
        run.font.bold = True

        col_w = content_w * 0.47
        gap = content_w * 0.06
        body_top = 1.8
        body_h = 4.5

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(body_top), Inches(col_w), Inches(body_h)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.left_content
        apply_text_style(run, font_body, 14, text_primary)

        # Divider line (center vertical)
        center_x = margin_l + col_w + gap / 2
        connector = slide.shapes.add_connector(
            1,  # straight connector
            Inches(center_x), Inches(body_top),
            Inches(center_x), Inches(body_top + body_h)
        )
        connector.line.color.rgb = hex_to_rgb(divider_color)
        connector.line.width = Pt(1)

        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(margin_l + col_w + gap), Inches(body_top),
            Inches(col_w), Inches(body_h)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.right_content
        apply_text_style(run, font_body, 14, text_primary)

        self._add_footer(slide)

    # ── Table Slide ─────────────────────────────────────────────

    def _render_table(self, prim: TableSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        text_light = self.colors.get("text_light", "#FFFFFF")
        text_primary = self.colors.get("text_primary", "#1a1a2e")
        card_bg = self.colors.get("card_background", "#F5F5F7")
        font_title = self.fonts.get("title", "Arial")
        font_body = self.fonts.get("body", "Arial")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(0.5), Inches(content_w), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.title
        apply_text_style(run, font_title, 28, primary)
        run.font.bold = True

        # Table
        rows_count = len(prim.rows) + 1  # +1 for header
        cols_count = len(prim.headers) if prim.headers else 1

        table_shape = slide.shapes.add_table(
            rows_count, cols_count,
            Inches(margin_l), Inches(1.8),
            Inches(content_w), Inches(min(rows_count * 0.5, 5.0))
        )
        table = table_shape.table

        # Style header row
        for col_idx, header in enumerate(prim.headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            # Header cell styling
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = hex_to_rgb(primary)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    apply_text_style(run, font_body, 12, text_light)
                    run.font.bold = True
            # Padding
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

        # Data rows with zebra striping
        for row_idx, row_data in enumerate(prim.rows):
            is_even = row_idx % 2 == 0
            bg = card_bg if is_even else "#FFFFFF"
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= cols_count:
                    break
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = hex_to_rgb(bg)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        apply_text_style(run, font_body, 12, text_primary)
                cell.margin_left = Inches(0.1)
                cell.margin_right = Inches(0.1)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)

        self._add_footer(slide)

    # ── Timeline Slide ──────────────────────────────────────────

    def _render_timeline(self, prim: TimelineSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        accent = self.colors.get("accent", "#0f3460")
        secondary = self.colors.get("secondary", "#16213e")
        text_primary = self.colors.get("text_primary", "#1a1a2e")
        font_title = self.fonts.get("title", "Arial")
        font_body = self.fonts.get("body", "Arial")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(0.5), Inches(content_w), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.title
        apply_text_style(run, font_title, 28, primary)
        run.font.bold = True

        events = prim.events
        n = len(events)
        if n == 0:
            self._add_footer(slide)
            return

        # Horizontal timeline
        line_y = 3.5
        line_left = margin_l + 0.5
        line_right = self.slide_w - margin_l - 0.5
        line_w = line_right - line_left

        # Connection line
        connector = slide.shapes.add_connector(
            1, Inches(line_left), Inches(line_y),
            Inches(line_right), Inches(line_y)
        )
        connector.line.color.rgb = hex_to_rgb(secondary)
        connector.line.width = Pt(2)

        # Event markers and labels
        for i, event in enumerate(events):
            if n == 1:
                x = line_left + line_w / 2
            else:
                x = line_left + (i / (n - 1)) * line_w

            # Circle marker
            marker_size = 0.25
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - marker_size / 2), Inches(line_y - marker_size / 2),
                Inches(marker_size), Inches(marker_size)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = hex_to_rgb(accent)
            marker.line.fill.background()

            # Date label — above the line
            date_box = slide.shapes.add_textbox(
                Inches(x - 1), Inches(line_y - 1.2),
                Inches(2), Inches(0.8)
            )
            tf = date_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = event.get("date", "")
            apply_text_style(run, font_body, 12, primary)
            run.font.bold = True

            # Description — below the line
            desc_box = slide.shapes.add_textbox(
                Inches(x - 1), Inches(line_y + 0.4),
                Inches(2), Inches(1.5)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = event.get("description", "")
            apply_text_style(run, font_body, 11, text_primary)

        self._add_footer(slide)

    # ── Chart Slide ─────────────────────────────────────────────

    def _render_chart(self, prim: ChartSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        font_title = self.fonts.get("title", "Arial")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(0.5), Inches(content_w), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = prim.title
        apply_text_style(run, font_title, 28, primary)
        run.font.bold = True

        chart_colors = self.config.get("chart_colors", [])
        add_chart_to_slide(
            slide, prim.chart_type, prim.data, prim.labels,
            Inches(margin_l), Inches(1.8), Inches(content_w), Inches(5),
            chart_colors=chart_colors
        )
        self._add_footer(slide)

    # ── Blank Slide ─────────────────────────────────────────────

    def _render_blank(self, prim: BlankSlide):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._add_footer(slide)

    # ── Fallback ────────────────────────────────────────────────

    def _render_fallback(self, prim: SlidePrimitive):
        slide = self.prs.slides.add_slide(self.blank_layout)
        primary = self.colors.get("primary", "#1a1a2e")
        font_title = self.fonts.get("title", "Arial")
        margin_l = self.spacing.get("margin_left", 0.8)
        content_w = self.slide_w - margin_l * 2

        title_box = slide.shapes.add_textbox(
            Inches(margin_l), Inches(0.5), Inches(content_w), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = getattr(prim, 'title', 'Untitled Slide')
        apply_text_style(run, font_title, 28, primary)
        run.font.bold = True
        self._add_footer(slide)
